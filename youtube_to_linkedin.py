#!/usr/bin/env python3
"""
Pipeline: Latest YouTube video → transcript → 5-slide PPT/PDF + LinkedIn post

Usage:
  python3 youtube_to_linkedin.py                          # searches "Limitless AI"
  python3 youtube_to_linkedin.py @LimitlessAI             # specific handle
  python3 youtube_to_linkedin.py "https://youtu.be/ID"   # direct video URL
  python3 youtube_to_linkedin.py "search:my query"       # explicit search
"""

import sys
import re
import textwrap
import ssl
import urllib.request
import xml.etree.ElementTree as ET
from pathlib import Path

# SSL bypass for environments with intercepting proxies
ssl._create_default_https_context = ssl._create_unverified_context

import yt_dlp
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from reportlab.lib.enums import TA_CENTER

# ── Config ──────────────────────────────────────────────────────────────────
INPUT   = sys.argv[1] if len(sys.argv) > 1 else "search:Limitless AI latest"
OUT_DIR = Path("output")
OUT_DIR.mkdir(exist_ok=True)

# Brand colours
BG_DARK = RGBColor(0x08, 0x09, 0x0d)
GOLD    = RGBColor(0xC9, 0xA8, 0x4C)
WHITE   = RGBColor(0xF0, 0xF1, 0xF5)
GREY    = RGBColor(0x66, 0x68, 0x78)

YDL_BASE = {
    "quiet": True,
    "skip_download": True,
    "nocheckcertificate": True,
}


# ── Step 1: resolve input → video metadata ──────────────────────────────────
def _ydl_info(url: str, extra: dict = {}) -> dict | None:
    opts = {**YDL_BASE, **extra}
    try:
        with yt_dlp.YoutubeDL(opts) as ydl:
            return ydl.extract_info(url, download=False)
    except Exception:
        return None


def _entry_to_video(entry: dict) -> dict:
    vid_id = entry.get("id") or entry.get("url", "").split("v=")[-1]
    return {
        "id":    vid_id,
        "title": entry.get("title", "Untitled"),
        "url":   f"https://www.youtube.com/watch?v={vid_id}",
    }


def _channel_rss(channel_id: str) -> dict | None:
    """Try YouTube RSS feed as a lightweight fallback."""
    feed_url = f"https://www.youtube.com/feeds/videos.xml?channel_id={channel_id}"
    try:
        with urllib.request.urlopen(feed_url, timeout=10) as r:
            root = ET.fromstring(r.read())
        ns = {
            "atom": "http://www.w3.org/2005/Atom",
            "yt":   "http://www.youtube.com/xml/schemas/2015",
        }
        entry = root.findall("atom:entry", ns)[0]
        vid_id = entry.find("yt:videoId", ns).text
        title  = entry.find("atom:title", ns).text
        return {"id": vid_id, "title": title, "url": f"https://www.youtube.com/watch?v={vid_id}"}
    except Exception:
        return None


def get_latest_video(input_str: str) -> dict:
    flat_opts = {"extract_flat": "in_playlist", "playlist_items": "1"}

    # 1. Direct video URL
    if "watch?v=" in input_str or "youtu.be/" in input_str:
        info = _ydl_info(input_str)
        if info:
            return {"id": info["id"], "title": info.get("title", "Untitled"),
                    "url": f"https://www.youtube.com/watch?v={info['id']}"}

    # 2. Explicit search: prefix
    if input_str.startswith("search:"):
        query = input_str[7:]
        info = _ydl_info(f"ytsearch1:{query}", flat_opts)
        if info and info.get("entries"):
            return _entry_to_video(list(info["entries"])[0])

    # 3. Channel URL / handle
    base = input_str.rstrip("/")
    if not base.startswith("http"):
        base = "https://www.youtube.com/" + base.lstrip("@")
        if not base.startswith("https://www.youtube.com/@"):
            base = "https://www.youtube.com/@" + input_str.lstrip("@")

    for suffix in ["/videos", "/streams", ""]:
        info = _ydl_info(base + suffix, flat_opts)
        if info:
            entries = list(info.get("entries") or [])
            if entries:
                return _entry_to_video(entries[0])
            # No entries, but we have channel_id — try RSS
            cid = info.get("channel_id")
            if cid:
                rss = _channel_rss(cid)
                if rss:
                    return rss

    # 4. Last resort: YouTube search
    query = re.sub(r"https?://[^\s]+", "", input_str).strip() or input_str
    info = _ydl_info(f"ytsearch1:{query}", flat_opts)
    if info and info.get("entries"):
        return _entry_to_video(list(info["entries"])[0])

    raise RuntimeError(f"Could not resolve a video from: {input_str}")


# ── Step 2: fetch transcript ─────────────────────────────────────────────────
def get_transcript(video_id: str) -> str:
    try:
        api = YouTubeTranscriptApi()
        fetched = api.fetch(video_id)
        return " ".join(s.text for s in fetched)
    except Exception as e:
        print(f"[!] Transcript unavailable ({e}), using description fallback.")
        # Fetch video description via yt-dlp
        info = _ydl_info(f"https://www.youtube.com/watch?v={video_id}")
        desc = (info or {}).get("description", "")
        if desc:
            return desc
        return (
            "Full transcript not available for this video. "
            "Please watch the original episode on YouTube. " * 15
        )


# ── Step 3: split transcript into 5 slide sections ───────────────────────────
def parse_into_slides(title: str, transcript: str) -> list[dict]:
    words = transcript.split()
    total = len(words)
    chunk = max(total // 4, 1)
    parts = [" ".join(words[i * chunk:(i + 1) * chunk]) for i in range(4)]
    parts[3] = " ".join(words[3 * chunk:])

    headings = [
        "Introduction & Context",
        "Key Insights — Part 1",
        "Key Insights — Part 2",
        "Takeaways & What's Next",
    ]

    slides = [{"heading": title, "bullets": [], "is_title": True}]
    for heading, text in zip(headings, parts):
        sentences = re.split(r"(?<=[.!?])\s+", text.strip())
        bullets = []
        for sent in sentences:
            sent = sent.strip()
            if len(sent) > 20:
                bullets.append(sent[:200] + ("…" if len(sent) > 200 else ""))
            if len(bullets) == 5:
                break
        slides.append({"heading": heading, "bullets": bullets, "is_title": False})
    return slides  # 5 slides total


# ── Step 4a: build PowerPoint ────────────────────────────────────────────────
def _txt(slide, text, left, top, width, height,
         size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def build_pptx(slides: list[dict], video_url: str, out_path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_DARK

        W, H, pad = prs.slide_width, prs.slide_height, Inches(0.6)

        if slide_data["is_title"]:
            bar = slide.shapes.add_shape(1, 0, Inches(2.8), W, Inches(0.06))
            bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()

            _txt(slide, "LIMITLESS — Latest Episode",
                 pad, Inches(2.0), W - pad*2, Inches(0.5),
                 size=13, color=GOLD, align=PP_ALIGN.CENTER)
            _txt(slide, slide_data["heading"],
                 pad, Inches(3.0), W - pad*2, Inches(2.0),
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _txt(slide, video_url,
                 pad, Inches(5.5), W - pad*2, Inches(0.4),
                 size=10, color=GREY, align=PP_ALIGN.CENTER)
        else:
            bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.07))
            bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()

            _txt(slide, slide_data["heading"],
                 pad, Inches(0.3), W - pad*2, Inches(0.8),
                 size=28, bold=True, color=GOLD)

            for i, bullet in enumerate(slide_data["bullets"]):
                wrapped = textwrap.shorten(bullet, width=120, placeholder="…")
                _txt(slide, f"• {wrapped}",
                     pad + Inches(0.2), Inches(1.3) + i * Inches(0.8),
                     W - pad*2 - Inches(0.2), Inches(0.8),
                     size=15, color=WHITE)

            _txt(slide, str(idx + 1),
                 W - Inches(1), H - Inches(0.5), Inches(0.8), Inches(0.4),
                 size=11, color=GREY, align=PP_ALIGN.RIGHT)

    prs.save(str(out_path))
    print(f"[✓] PPT saved → {out_path}")


# ── Step 4b: build PDF ───────────────────────────────────────────────────────
def build_pdf(slides: list[dict], video_url: str, out_path: Path):
    doc = SimpleDocTemplate(str(out_path), pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    gold_hex = "#C9A84C"

    title_s = ParagraphStyle("ST", parent=styles["Title"],
                              fontSize=22, textColor=colors.HexColor(gold_hex), spaceAfter=6)
    head_s  = ParagraphStyle("SH", parent=styles["Heading1"],
                              fontSize=16, textColor=colors.HexColor(gold_hex), spaceAfter=4)
    body_s  = ParagraphStyle("SB", parent=styles["Normal"],
                              fontSize=11, leading=16, spaceAfter=4)
    meta_s  = ParagraphStyle("SM", parent=styles["Normal"],
                              fontSize=9, textColor=colors.grey, alignment=TA_CENTER)

    story = []
    for slide in slides:
        if slide["is_title"]:
            story += [Spacer(1, inch*0.5),
                      Paragraph("LIMITLESS — Latest Episode", meta_s),
                      Spacer(1, 10),
                      Paragraph(slide["heading"], title_s),
                      Spacer(1, 6),
                      Paragraph(video_url, meta_s)]
        else:
            story += [HRFlowable(width="100%", thickness=1,
                                 color=colors.HexColor(gold_hex)),
                      Spacer(1, 6),
                      Paragraph(slide["heading"], head_s)]
            for b in slide["bullets"]:
                story.append(Paragraph(f"• {b}", body_s))
            story.append(Spacer(1, 14))

    doc.build(story)
    print(f"[✓] PDF saved → {out_path}")


# ── Step 5: generate LinkedIn post ──────────────────────────────────────────
def build_linkedin_post(title: str, slides: list[dict], video_url: str) -> str:
    bullets = []
    for s in slides[1:]:
        bullets.extend(s["bullets"][:2])
    highlights = "\n".join(f"• {b}" for b in bullets[:6])

    return f"""\
Just watched the latest from Limitless 🎙️

📺 "{title}"

Here are my top takeaways:

{highlights}

Full episode → {video_url}

What do you think? Drop your thoughts below 👇

#Limitless #AI #Productivity #Innovation #FutureOfWork
""".strip()


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    print(f"[→] Resolving: {INPUT}")
    video = get_latest_video(INPUT)
    print(f"[✓] Video : {video['title']}")
    print(f"    URL   : {video['url']}")

    print("[→] Fetching transcript…")
    transcript = get_transcript(video["id"])
    print(f"[✓] Transcript: {len(transcript.split())} words")

    slides = parse_into_slides(video["title"], transcript)

    safe  = re.sub(r"[^\w\-]", "_", video["title"])[:60]
    pptx  = OUT_DIR / f"{safe}.pptx"
    pdf   = OUT_DIR / f"{safe}.pdf"
    post  = OUT_DIR / f"{safe}_linkedin_post.txt"

    build_pptx(slides, video["url"], pptx)
    build_pdf(slides, video["url"], pdf)

    linkedin = build_linkedin_post(video["title"], slides, video["url"])
    post.write_text(linkedin)
    print(f"[✓] LinkedIn post → {post}")

    print("\n" + "=" * 60)
    print("LINKEDIN POST PREVIEW")
    print("=" * 60)
    print(linkedin)
    print("=" * 60)


if __name__ == "__main__":
    main()
